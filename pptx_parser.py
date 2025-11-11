"""
PowerPoint speaker notes extraction and splitting logic
"""

import re
from pptx import Presentation
from typing import List, Tuple


def extract_slides(pptx_file) -> List[dict]:
    """
    Extract slide information from PowerPoint file.

    Returns list of dicts with:
    - slide_num: int
    - title: str
    - notes: str (full speaker notes text)
    """
    prs = Presentation(pptx_file)
    slides = []

    for i, slide in enumerate(prs.slides, 1):
        # Get slide title
        title = "Slide " + str(i)
        if slide.shapes.title and slide.shapes.title.text.strip():
            title = slide.shapes.title.text.strip()

        # Get speaker notes
        notes = ""
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame:
                notes = notes_slide.notes_text_frame.text.strip()

        slides.append({
            'slide_num': i,
            'title': title,
            'notes': notes
        })

    return slides


def split_notes(notes: str, level: int) -> List[str]:
    """
    Split speaker notes into chunks based on split level.

    Args:
        notes: Full speaker notes text
        level: 1-4 split granularity
            1: Whole note (no split)
            2: Double line breaks (paragraphs)
            3: Single line breaks
            4: Sentences (period + space)

    Returns:
        List of text chunks
    """
    if not notes or not notes.strip():
        return []

    notes = notes.strip()

    if level == 1:
        # Whole note
        return [notes]

    elif level == 2:
        # Split by double line breaks (paragraphs)
        chunks = re.split(r'\n\s*\n', notes)
        chunks = [c.strip() for c in chunks if c.strip()]
        return chunks if chunks else [notes]

    elif level == 3:
        # Split by single line breaks
        chunks = notes.split('\n')
        chunks = [c.strip() for c in chunks if c.strip()]
        return chunks if chunks else [notes]

    elif level == 4:
        # Split by sentences (period + space/newline)
        chunks = re.split(r'\.(?:\s+|\n+)', notes)
        chunks = [c.strip() + '.' if c.strip() and not c.strip().endswith('.') else c.strip()
                  for c in chunks if c.strip()]
        return chunks if chunks else [notes]

    else:
        # Default to whole note for any other value
        return [notes]


def create_button_label(title: str, chunk_index: int, total_chunks: int,
                        slide_num: int, content: str,
                        format_type: str = "num_part_content", max_length: int = 30) -> str:
    """
    Create button label with various format options.

    Args:
        title: Slide title
        chunk_index: Current chunk number (0-based)
        total_chunks: Total number of chunks for this slide
        slide_num: Slide number (1-based)
        content: The actual message content/chunk text
        format_type: Label format type:
            - "title_part": "Title (Part N)" [default]
            - "slide_content": "Slide N: Content..."
            - "content_only": "Content..."
            - "num_title": "N - Title"
            - "num_part_content": "N.P: Content..."
        max_length: Maximum label length before adding ellipsis

    Returns:
        Formatted button label
    """
    if format_type == "title_part":
        # "Title (Part N)" format
        truncated_title = title[:max_length-3] + "..." if len(title) > max_length else title
        if total_chunks > 1:
            label = f"{truncated_title} ({chunk_index + 1})"
        else:
            label = truncated_title

    elif format_type == "slide_content":
        # "Slide N: Content..." format
        content_preview = content[:max_length-10] + "..." if len(content) > max_length-10 else content
        # Replace newlines with spaces for cleaner preview
        content_preview = content_preview.replace('\n', ' ').strip()
        label = f"Slide {slide_num}: {content_preview}"

    elif format_type == "content_only":
        # "Content..." format
        content_preview = content[:max_length] + "..." if len(content) > max_length else content
        # Replace newlines with spaces for cleaner preview
        label = content_preview.replace('\n', ' ').strip()

    elif format_type == "num_title":
        # "N - Title" format
        truncated_title = title[:max_length-5] + "..." if len(title) > max_length-5 else title
        label = f"{slide_num} - {truncated_title}"

    elif format_type == "num_part_content":
        # "N.P: Content..." format
        if total_chunks > 1:
            part_prefix = f"{slide_num}.{chunk_index + 1}: "
        else:
            part_prefix = f"{slide_num}: "

        remaining_length = max_length - len(part_prefix)
        content_preview = content[:remaining_length] + "..." if len(content) > remaining_length else content
        # Replace newlines with spaces for cleaner preview
        content_preview = content_preview.replace('\n', ' ').strip()
        label = f"{part_prefix}{content_preview}"

    else:
        # Default to title_part format
        truncated_title = title[:max_length-3] + "..." if len(title) > max_length else title
        if total_chunks > 1:
            label = f"{truncated_title} ({chunk_index + 1})"
        else:
            label = truncated_title

    return label


def parse_pptx_to_buttons(pptx_file, split_levels: dict = None, default_level: int = 2,
                          label_format: str = "title_part", max_label_length: int = 30) -> List[Tuple[str, str, int]]:
    """
    Parse PowerPoint file into button data.

    Args:
        pptx_file: PowerPoint file object
        split_levels: Dict mapping slide_num to split level (overrides default)
        default_level: Default split level for all slides
        label_format: Label format type (title_part, slide_content, content_only, num_title, num_part_content)
        max_label_length: Maximum label length before truncation

    Returns:
        List of tuples: (label, message, slide_num)
    """
    slides = extract_slides(pptx_file)
    buttons = []

    for slide in slides:
        slide_num = slide['slide_num']
        title = slide['title']
        notes = slide['notes']

        # Skip slides with no notes
        if not notes:
            continue

        # Determine split level for this slide
        level = split_levels.get(slide_num, default_level) if split_levels else default_level

        # Split notes
        chunks = split_notes(notes, level)

        # Create buttons
        for i, chunk in enumerate(chunks):
            label = create_button_label(
                title=title,
                chunk_index=i,
                total_chunks=len(chunks),
                slide_num=slide_num,
                content=chunk,
                format_type=label_format,
                max_length=max_label_length
            )
            buttons.append((label, chunk, slide_num))

    return buttons
